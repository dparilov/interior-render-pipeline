#!/usr/bin/env python3
"""
Чистая презентация без шаблона — 6 слайдов
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_PATH = '/home/dima/.openclaw/workspace/Metodika_clean.pptx'

# Цвета T1
T1_BLUE = RGBColor(0, 120, 215)
T1_DARK = RGBColor(30, 30, 30)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(100, 100, 100)

def add_title_slide(prs, title, subtitle):
    """Титульный слайд"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = T1_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Подзаголовок
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    """Контентный слайд"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = T1_BLUE
    
    # Линия под заголовком
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = T1_BLUE
    line.line.fill.background()
    
    # Контент
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line_text in enumerate(content_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(16)
        p.font.color.rgb = T1_DARK
        
        # Секции жирным
        if line_text.endswith(":") or line_text.startswith("→"):
            p.font.bold = True
            p.font.size = Pt(17)
        if line_text.startswith("•"):
            p.font.size = Pt(15)
        if line_text == "":
            p.font.size = Pt(10)
    
    return slide

def add_final_slide(prs, text="Спасибо за внимание"):
    """Финальный слайд"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.3), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = T1_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # === СЛАЙД 1: Титульный ===
    add_title_slide(prs,
        "Методика оценки эффективности\nспонсорских мероприятий",
        "Управление по работе с бизнес-мероприятиями | Холдинг Т1 | 2025"
    )
    
    # === СЛАЙД 2: Проблема + Решение ===
    add_content_slide(prs, "Проблема и решение", [
        "ПРОБЛЕМА:",
        "• Нет единой метрики для сравнения мероприятий разных форматов",
        "• Субъективность оценки спонсорских пакетов",
        "• Сложно обосновать ROI спонсорства перед руководством",
        "• Нет системы для решений «участвовать / не участвовать»",
        "",
        "РЕШЕНИЕ — Матрица оценки эффективности:",
        "• Количественные метрики (охват, лиды, конверсия, ROI)",
        "• Качественная оценка спонсорских опций (коэффициенты k от 1 до 5)",
        "• Комплексный балл мероприятия (КБМ) для сравнения",
        "",
        "→ Объективное сравнение любых мероприятий по единой шкале"
    ])
    
    # === СЛАЙД 3: Архитектура ===
    add_content_slide(prs, "Как устроена матрица", [
        "ВВОД ДАННЫХ (ивент-менеджер по данным от организаторов):",
        "• Общая информация: название, дата, город, участник, приоритет",
        "• Исходные данные: участники, встречи, контакты, бюджет, доход, NPS",
        "• Оценка 8 спонсорских опций (k от 1 до 5)",
        "",
        "                                              ↓",
        "",
        "АВТОМАТИЧЕСКИЙ РАСЧЁТ:",
        "• Метрики: OTS, GRP, CPL, CPA, ROI, ROMI, NPS",
        "• Индекс спонсорских опций (Idx) = среднее k по всем опциям",
        "• Комплексный балл мероприятия (КБМ)",
        "",
        "                                              ↓",
        "",
        "→ РЕЗУЛЬТАТ: рейтинг мероприятий для сравнения и планирования"
    ])
    
    # === СЛАЙД 4: КБМ + Пример ===
    add_content_slide(prs, "Комплексный балл и пример расчёта", [
        "ФОРМУЛА КБМ:",
        "КБМ = w₁×GRP + w₂×Idx + w₃×(1/CPL) + w₄×CR",
        "",
        "Веса: GRP 25% | Idx 20% | 1/CPL 30% | CR 25% (настраиваются под цели)",
        "",
        "ПРИМЕР — мероприятия 2024-2025:",
        "",
        "Мероприятие          OTS         GRP         Idx",
        "─────────────────────────────────────────────────",
        "ПМГФ 2024            34 400      0,54%       2,75",
        "Smart Mining         554         15,16%      3,33",
        "Белые ночи           375         20,27%      2,75",
        "",
        "→ Smart Mining и Белые ночи: выше вовлечённость при меньшем охвате",
        "   — эффективнее для B2B-лидогенерации"
    ])
    
    # === СЛАЙД 5: Выводы ===
    add_content_slide(prs, "Выводы и следующие шаги", [
        "ЧТО ДАЁТ МЕТОДИКА:",
        "• Объективное сравнение любых мероприятий",
        "• Обоснование бюджетов на спонсорство",
        "• База для решений «участвовать / не участвовать»",
        "",
        "ПРОЦЕСС ИСПОЛЬЗОВАНИЯ:",
        "• Ивент-менеджер заполняет матрицу по данным от организаторов",
        "• Метрики и КБМ рассчитываются автоматически",
        "• Формируется рейтинг для планирования на следующий год",
        "",
        "СЛЕДУЮЩИЕ ШАГИ:",
        "• Утвердить веса KPI под приоритеты Т1",
        "• Пилот на мероприятиях Q2 2025",
        "• Сформировать рейтинг по итогам года"
    ])
    
    # === СЛАЙД 6: Спасибо ===
    add_final_slide(prs, "Спасибо за внимание")
    
    prs.save(OUTPUT_PATH)
    print(f"✅ Сохранено: {OUTPUT_PATH}")
    print(f"   Слайдов: {len(prs.slides)}")

if __name__ == '__main__':
    create_presentation()

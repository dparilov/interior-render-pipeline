#!/usr/bin/env python3
"""
Создание презентации "Методика оценки эффективности спонсорских мероприятий"
на основе шаблона T1
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy

# Пути к файлам
TEMPLATE_PATH = '/home/dima/.openclaw/media/inbound/RUS_Template_T1_окт24_облегченная_версия_2---e96271ef-6531-4673-9d93-5559d2895d14.pptx'
OUTPUT_PATH = '/home/dima/.openclaw/workspace/Методика_оценки_спонсорских_мероприятий_T1.pptx'

def create_presentation():
    # Загружаем шаблон
    prs = Presentation(TEMPLATE_PATH)
    
    # Получаем макеты
    layouts = {layout.name: layout for layout in prs.slide_layouts}
    
    # Удаляем все слайды из шаблона
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # === СЛАЙД 1: Титульный ===
    slide1 = prs.slides.add_slide(layouts['Титульный слайд 1'])
    for shape in slide1.shapes:
        if hasattr(shape, 'text_frame'):
            if 'Шаблон' in shape.text or 'презентаций' in shape.text.lower():
                shape.text_frame.paragraphs[0].runs[0].text = "Методика оценки эффективности\nспонсорских мероприятий"
    
    # === СЛАЙД 2: Проблема (разделитель) ===
    slide2 = prs.slides.add_slide(layouts['Разделитель голубой'])
    for shape in slide2.shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if 'Обложка' in run.text or 'раздел' in run.text.lower():
                        run.text = "Проблема"
    
    # === СЛАЙД 3: Описание проблемы ===
    slide3 = prs.slides.add_slide(layouts['Контентный слайд'])
    for shape in slide3.shapes:
        if hasattr(shape, 'text_frame'):
            if 'Заголовок' in shape.text:
                shape.text_frame.paragraphs[0].runs[0].text = "Почему сложно оценить отдачу от спонсорства"
            elif shape.text_frame.paragraphs:
                tf = shape.text_frame
                tf.clear()
                
                problems = [
                    "Нет единой методики оценки — каждое мероприятие считается по-своему",
                    "Субъективность — «было круто» vs измеримые результаты", 
                    "Сложно сравнивать — ПМГФ vs камерный бизнес-завтрак",
                    "Разные спонсорские пакеты — как сравнить стенд и сессию?",
                    "Отложенный эффект — сделки закрываются через 3-6 месяцев"
                ]
                
                for i, problem in enumerate(problems):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"• {problem}"
                    p.level = 0
    
    # === СЛАЙД 4: Решение (разделитель) ===
    slide4 = prs.slides.add_slide(layouts['Разделитель темный'])
    for shape in slide4.shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if 'Обложка' in run.text or 'раздел' in run.text.lower():
                        run.text = "Решение"
    
    # === СЛАЙД 5: Комплексный индекс ===
    slide5 = prs.slides.add_slide(layouts['Контентный слайд'])
    for shape in slide5.shapes:
        if hasattr(shape, 'text_frame'):
            if 'Заголовок' in shape.text:
                shape.text_frame.paragraphs[0].runs[0].text = "Комплексный индекс оценки мероприятий"
            elif shape.text_frame.paragraphs:
                tf = shape.text_frame
                tf.clear()
                
                points = [
                    "Единая матрица для всех типов мероприятий",
                    "Объективные метрики: OTS, GRP, CPL, ROI, NPS",
                    "Оценка качества спонсорского пакета (Idx)",
                    "Комплексный балл мероприятия (КБМ) — итоговый рейтинг",
                    "Возможность сравнивать и ранжировать мероприятия"
                ]
                
                for i, point in enumerate(points):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"• {point}"
                    p.level = 0
    
    # === СЛАЙД 6: Архитектура матрицы ===
    slide6 = prs.slides.add_slide(layouts['Контентный слайд'])
    for shape in slide6.shapes:
        if hasattr(shape, 'text_frame'):
            if 'Заголовок' in shape.text:
                shape.text_frame.paragraphs[0].runs[0].text = "Архитектура матрицы: 5 блоков"
            elif shape.text_frame.paragraphs:
                tf = shape.text_frame
                tf.clear()
                
                blocks = [
                    "Блок 1. Общая информация — название, дата, город, участник, приоритет",
                    "Блок 2. Исходные данные (ввод) — участники, встречи, контакты, бюджет, доход",
                    "Блок 3. Расчётные метрики (авто) — OTS, GRP, CPL, CPA, ROI, NPS и др.",
                    "Блок 4. Индекс спонсорских опций — оценка 8 параметров (k от 1 до 5)",
                    "Блок 5. Комплексный балл мероприятия (КБМ) — итоговый рейтинг"
                ]
                
                for i, block in enumerate(blocks):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = block
                    p.level = 0
    
    # === СЛАЙД 7: Ключевые метрики ===
    slide7 = prs.slides.add_slide(layouts['Контентный слайд'])
    for shape in slide7.shapes:
        if hasattr(shape, 'text_frame'):
            if 'Заголовок' in shape.text:
                shape.text_frame.paragraphs[0].runs[0].text = "Ключевые метрики"
            elif shape.text_frame.paragraphs:
                tf = shape.text_frame
                tf.clear()
                
                metrics = [
                    "OTS (Opportunity To See) — максимальный охват аудитории",
                    "GRP (Gross Rating Point) — доля аудитории с прямым контактом",
                    "CPL (Cost Per Lead) — стоимость одного лида",
                    "CPA (Cost Per Action) — стоимость целевого действия",
                    "ROI / ROMI — возврат инвестиций",
                    "NPS (Net Promoter Score) — индекс лояльности",
                    "CR (Conversion Rate) — конверсия лид → сделка"
                ]
                
                for i, metric in enumerate(metrics):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"• {metric}"
                    p.level = 0
    
    # === СЛАЙД 8: Индекс спонсорских опций ===
    slide8 = prs.slides.add_slide(layouts['Контентный слайд'])
    for shape in slide8.shapes:
        if hasattr(shape, 'text_frame'):
            if 'Заголовок' in shape.text:
                shape.text_frame.paragraphs[0].runs[0].text = "Индекс спонсорских опций (Idx)"
            elif shape.text_frame.paragraphs:
                tf = shape.text_frame
                tf.clear()
                
                content = [
                    "8 параметров оценки (p1–p8):",
                    "",
                    "p1: Выставочный стенд",
                    "p2: Сессия в деловой программе", 
                    "p3: Размещение на сайте мероприятия",
                    "p4: Логотип в программе / на баннерах",
                    "p5: Сувениры / мерч в сумке участника",
                    "p6: Активации / игровые механики",
                    "p7: Нетворкинг / бизнес-завтрак",
                    "p8: Прочее",
                    "",
                    "Каждый параметр оценивается k от 1 до 5",
                    "Idx = Σ(k₁...kₙ) / n — средний балл по всем опциям"
                ]
                
                for i, line in enumerate(content):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    p.level = 0
    
    # === СЛАЙД 9: Матрица коэффициентов k ===
    slide9 = prs.slides.add_slide(layouts['Контентный слайд'])
    for shape in slide9.shapes:
        if hasattr(shape, 'text_frame'):
            if 'Заголовок' in shape.text:
                shape.text_frame.paragraphs[0].runs[0].text = "Матрица коэффициентов значимости (k)"
            elif shape.text_frame.paragraphs:
                tf = shape.text_frame
                tf.clear()
                
                content = [
                    "Пример: Выставочный стенд",
                    "",
                    "k=5 — Центральный, у входа, ≥70 кв.м.",
                    "k=4 — Ключевой павильон, 1-я линия",
                    "k=3 — 2-я линия, видимый",
                    "k=2 — Угловой, малая площадь",
                    "k=1 — Задворки, плохая проходимость",
                    "",
                    "Аналогичная шкала для всех 8 опций",
                    "(подробная матрица — в справочнике)"
                ]
                
                for i, line in enumerate(content):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    p.level = 0
    
    # === СЛАЙД 10: Комплексный балл (КБМ) ===
    slide10 = prs.slides.add_slide(layouts['Контентный слайд'])
    for shape in slide10.shapes:
        if hasattr(shape, 'text_frame'):
            if 'Заголовок' in shape.text:
                shape.text_frame.paragraphs[0].runs[0].text = "Комплексный балл мероприятия (КБМ)"
            elif shape.text_frame.paragraphs:
                tf = shape.text_frame
                tf.clear()
                
                content = [
                    "Формула:",
                    "КБМ = w1×GRP + w2×Idx + w3×(1/CPL) + w4×NPS + w5×ROI",
                    "",
                    "Рекомендуемые веса:",
                    "• w1 (GRP) = 25% — охват и вовлечённость",
                    "• w2 (Idx) = 20% — качество спонсорского пакета",
                    "• w3 (1/CPL) = 30% — эффективность лидогенерации",
                    "• w4 (NPS) = 15% — удовлетворённость",
                    "• w5 (ROI) = 10% — финансовая отдача",
                    "",
                    "Веса настраиваются под приоритеты компании"
                ]
                
                for i, line in enumerate(content):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    p.level = 0
    
    # === СЛАЙД 11: Пример расчёта (разделитель) ===
    slide11 = prs.slides.add_slide(layouts['Разделитель голубой'])
    for shape in slide11.shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if 'Обложка' in run.text or 'раздел' in run.text.lower():
                        run.text = "Пример расчёта"
    
    # === СЛАЙД 12: Пример на реальных мероприятиях ===
    slide12 = prs.slides.add_slide(layouts['Контентный слайд'])
    for shape in slide12.shapes:
        if hasattr(shape, 'text_frame'):
            if 'Заголовок' in shape.text:
                shape.text_frame.paragraphs[0].runs[0].text = "Сравнение мероприятий 2024-2025"
            elif shape.text_frame.paragraphs:
                tf = shape.text_frame
                tf.clear()
                
                content = [
                    "                          ПМГФ        Smart Mining    Белые ночи",
                    "",
                    "OTS (охват)              34 400            554              375",
                    "Касания                     186              84               76",
                    "GRP                       0,54%         15,16%          20,27%",
                    "Idx (опции)                2,75           3,33            2,75",
                    "",
                    "Вывод: Smart Mining и Белые ночи показали",
                    "более высокую вовлечённость при меньшем охвате"
                ]
                
                for i, line in enumerate(content):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    p.level = 0
    
    # === СЛАЙД 13: Выводы и next steps ===
    slide13 = prs.slides.add_slide(layouts['Контентный слайд'])
    for shape in slide13.shapes:
        if hasattr(shape, 'text_frame'):
            if 'Заголовок' in shape.text:
                shape.text_frame.paragraphs[0].runs[0].text = "Выводы и следующие шаги"
            elif shape.text_frame.paragraphs:
                tf = shape.text_frame
                tf.clear()
                
                content = [
                    "Что даёт методика:",
                    "• Объективное сравнение любых мероприятий",
                    "• Обоснование бюджетов на спонсорство",
                    "• Выявление наиболее эффективных форматов",
                    "• База для планирования на следующий год",
                    "",
                    "Следующие шаги:",
                    "• Утвердить веса KPI под цели Т1",
                    "• Внедрить матрицу для всех мероприятий 2025",
                    "• Собрать данные и сформировать рейтинг"
                ]
                
                for i, line in enumerate(content):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    p.level = 0
    
    # === СЛАЙД 14: Спасибо ===
    slide14 = prs.slides.add_slide(layouts['1_Спасибо за внимание'])
    
    # Сохраняем
    prs.save(OUTPUT_PATH)
    print(f"✅ Презентация сохранена: {OUTPUT_PATH}")
    print(f"   Всего слайдов: {len(prs.slides)}")

if __name__ == '__main__':
    create_presentation()

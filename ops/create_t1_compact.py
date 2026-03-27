#!/usr/bin/env python3
"""
Компактная версия презентации — 6 слайдов для C-level
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

TEMPLATE_PATH = '/home/dima/.openclaw/media/inbound/RUS_Template_T1_окт24_облегченная_версия_2---e96271ef-6531-4673-9d93-5559d2895d14.pptx'
OUTPUT_PATH = '/home/dima/.openclaw/workspace/Metodika_T1_compact.pptx'

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=None):
    """Добавляет текстовый блок"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return shape

def add_block(slide, left, top, width, height, title, content, fill_color):
    """Добавляет цветной блок со скруглёнными углами"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    
    # Заголовок блока
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Содержимое
    p = tf.add_paragraph()
    p.text = content
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    return shape

def create_presentation():
    prs = Presentation(TEMPLATE_PATH)
    layouts = {layout.name: layout for layout in prs.slide_layouts}
    
    # Удаляем все слайды
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    content_layout = layouts.get('Контентный слайд')
    blank_layout = layouts.get('Пустой слайд')
    
    # === СЛАЙД 1: Титульный ===
    slide1 = prs.slides.add_slide(layouts['Титульный слайд 1'])
    for shape in slide1.placeholders:
        if shape.placeholder_format.type == 1:  # TITLE
            shape.text = "Методика оценки эффективности\nспонсорских мероприятий"
        elif shape.placeholder_format.type == 2:  # SUBTITLE
            shape.text = "Управление по работе с бизнес-мероприятиями\nХолдинг Т1 | 2025"
    
    # === СЛАЙД 2: Проблема + Решение ===
    slide2 = prs.slides.add_slide(content_layout)
    for shape in slide2.placeholders:
        if shape.placeholder_format.type == 1:
            shape.text = "Проблема и решение"
        elif shape.placeholder_format.type == 2:
            tf = shape.text_frame
            tf.clear()
            
            lines = [
                "ПРОБЛЕМА:",
                "• Нет единой метрики для сравнения мероприятий",
                "• Субъективность оценки спонсорских пакетов",
                "• Сложно обосновать ROI перед руководством",
                "",
                "РЕШЕНИЕ — Матрица оценки эффективности:",
                "• Количественные метрики (охват, лиды, ROI)",
                "• Качественная оценка опций (коэффициенты k)",
                "• Комплексный балл мероприятия (КБМ)",
                "",
                "→ Объективное сравнение любых мероприятий по единой шкале"
            ]
            
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                if line.startswith("ПРОБЛЕМА") or line.startswith("РЕШЕНИЕ"):
                    p.font.bold = True
    
    # === СЛАЙД 3: Архитектура (блоки) ===
    slide3 = prs.slides.add_slide(content_layout)
    for shape in slide3.placeholders:
        if shape.placeholder_format.type == 1:
            shape.text = "Как устроена матрица"
        elif shape.placeholder_format.type == 2:
            tf = shape.text_frame
            tf.clear()
            
            # Вместо сложной графики — структурированный текст
            lines = [
                "┌─────────────────────────────────────────────────────────┐",
                "│  ВВОД ДАННЫХ (ивент-менеджер)                          │",
                "│  • Участники, встречи, контакты, бюджет, доход         │",
                "│  • Оценка 8 спонсорских опций (k от 1 до 5)            │",
                "└─────────────────────────────────────────────────────────┘",
                "                           ↓",
                "┌─────────────────────────────────────────────────────────┐",
                "│  АВТОМАТИЧЕСКИЙ РАСЧЁТ                                 │",
                "│  • Метрики: OTS, GRP, CPL, CPA, ROI, NPS               │",
                "│  • Индекс опций (Idx) = среднее k                      │",
                "│  • Комплексный балл (КБМ)                              │",
                "└─────────────────────────────────────────────────────────┘",
                "                           ↓",
                "┌─────────────────────────────────────────────────────────┐",
                "│  РЕЗУЛЬТАТ: рейтинг мероприятий для сравнения          │",
                "└─────────────────────────────────────────────────────────┘"
            ]
            
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(11)
    
    # === СЛАЙД 4: КБМ + Пример ===
    slide4 = prs.slides.add_slide(content_layout)
    for shape in slide4.placeholders:
        if shape.placeholder_format.type == 1:
            shape.text = "Комплексный балл и пример"
        elif shape.placeholder_format.type == 2:
            tf = shape.text_frame
            tf.clear()
            
            lines = [
                "ФОРМУЛА КБМ:",
                "КБМ = 25%×GRP + 20%×Idx + 30%×(1/CPL) + 25%×CR",
                "",
                "ПРИМЕР — сравнение мероприятий 2024-2025:",
                "",
                "Мероприятие      | OTS      | GRP     | Idx  ",
                "─────────────────|──────────|─────────|──────",
                "ПМГФ 2024        | 34 400   | 0,54%   | 2,75 ",
                "Smart Mining     | 554      | 15,16%  | 3,33 ",
                "Белые ночи       | 375      | 20,27%  | 2,75 ",
                "",
                "→ Smart Mining и Белые ночи: выше вовлечённость (GRP)",
                "  при меньшем охвате — эффективнее для B2B"
            ]
            
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                if line.startswith("ФОРМУЛА") or line.startswith("ПРИМЕР"):
                    p.font.bold = True
                if "КБМ =" in line or line.startswith("→"):
                    p.font.bold = True
    
    # === СЛАЙД 5: Выводы и next steps ===
    slide5 = prs.slides.add_slide(content_layout)
    for shape in slide5.placeholders:
        if shape.placeholder_format.type == 1:
            shape.text = "Выводы и следующие шаги"
        elif shape.placeholder_format.type == 2:
            tf = shape.text_frame
            tf.clear()
            
            lines = [
                "ЧТО ДАЁТ МЕТОДИКА:",
                "• Объективное сравнение любых мероприятий",
                "• Обоснование бюджетов на спонсорство",
                "• База для решений «участвовать / не участвовать»",
                "",
                "ПРОЦЕСС ИСПОЛЬЗОВАНИЯ:",
                "• Ивент-менеджер заполняет матрицу по данным от организаторов",
                "• Метрики и КБМ рассчитываются автоматически",
                "• Формируется рейтинг для планирования на следующий год",
                "",
                "СЛЕДУЮЩИЕ ШАГИ:",
                "• Утвердить веса KPI под цели Т1",
                "• Пилот на мероприятиях Q2 2025",
                "• Сформировать рейтинг по итогам года"
            ]
            
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                if line.endswith(":"):
                    p.font.bold = True
    
    # === СЛАЙД 6: Спасибо ===
    slide6 = prs.slides.add_slide(layouts['1_Спасибо за внимание'])
    
    prs.save(OUTPUT_PATH)
    print(f"✅ Сохранено: {OUTPUT_PATH}")
    print(f"   Слайдов: {len(prs.slides)}")

if __name__ == '__main__':
    create_presentation()

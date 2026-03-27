#!/usr/bin/env python3
"""
Создание презентации v2 — с правильной работой с плейсхолдерами
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

TEMPLATE_PATH = '/home/dima/.openclaw/media/inbound/RUS_Template_T1_окт24_облегченная_версия_2---e96271ef-6531-4673-9d93-5559d2895d14.pptx'
OUTPUT_PATH = '/home/dima/.openclaw/workspace/Методика_оценки_T1_v2.pptx'

def set_text_in_placeholder(slide, ph_idx, text):
    """Устанавливает текст в плейсхолдер по индексу"""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            shape.text = text
            return True
    return False

def add_content_slide(prs, layout, title, bullets):
    """Добавляет слайд с заголовком и буллетами"""
    slide = prs.slides.add_slide(layout)
    
    # Ищем плейсхолдеры
    title_shape = None
    body_shape = None
    
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        # 1 = TITLE, 2 = BODY
        if ph_type == 1 or 'заголовок' in shape.name.lower() or 'title' in shape.name.lower():
            title_shape = shape
        elif ph_type == 2 or 'content' in shape.name.lower() or 'text' in shape.name.lower():
            body_shape = shape
    
    # Если не нашли по типу, берём по порядку
    placeholders = list(slide.placeholders)
    if not title_shape and len(placeholders) > 0:
        title_shape = placeholders[0]
    if not body_shape and len(placeholders) > 1:
        body_shape = placeholders[1]
    
    if title_shape:
        title_shape.text = title
    
    if body_shape and bullets:
        tf = body_shape.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
    
    return slide

def create_presentation():
    prs = Presentation(TEMPLATE_PATH)
    layouts = {layout.name: layout for layout in prs.slide_layouts}
    
    # Удаляем все слайды
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # Выбираем основной layout для контента
    content_layout = layouts.get('Контентный слайд') or layouts.get('Контент + подзаголовок')
    divider_blue = layouts.get('Разделитель голубой')
    divider_dark = layouts.get('Разделитель темный')
    
    # === СЛАЙД 1: Титульный ===
    slide1 = prs.slides.add_slide(layouts['Титульный слайд 1'])
    for shape in slide1.placeholders:
        if shape.placeholder_format.type == 1:  # TITLE
            shape.text = "Методика оценки эффективности\nспонсорских мероприятий"
        elif shape.placeholder_format.type == 2:  # SUBTITLE
            shape.text = "Холдинг Т1 | 2025"
    
    # === СЛАЙД 2: Разделитель "Проблема" ===
    slide2 = prs.slides.add_slide(divider_blue)
    for shape in slide2.placeholders:
        shape.text = "Проблема"
    
    # === СЛАЙД 3: Проблема ===
    add_content_slide(prs, content_layout,
        "Почему сложно оценить отдачу от спонсорства",
        [
            "• Нет единой методики — каждое мероприятие считается по-своему",
            "• Субъективность — «было круто» vs измеримые результаты",
            "• Сложно сравнивать — ПМГФ vs камерный бизнес-завтрак",
            "• Разные пакеты — как сравнить стенд и сессию?",
            "• Отложенный эффект — сделки через 3-6 месяцев"
        ])
    
    # === СЛАЙД 4: Разделитель "Решение" ===
    slide4 = prs.slides.add_slide(divider_dark)
    for shape in slide4.placeholders:
        shape.text = "Решение"
    
    # === СЛАЙД 5: Комплексный индекс ===
    add_content_slide(prs, content_layout,
        "Комплексный индекс оценки мероприятий",
        [
            "• Единая матрица для всех типов мероприятий",
            "• Объективные метрики: OTS, GRP, CPL, ROI, NPS",
            "• Оценка качества спонсорского пакета (Idx)",
            "• Комплексный балл мероприятия (КБМ) — итоговый рейтинг",
            "• Возможность сравнивать и ранжировать"
        ])
    
    # === СЛАЙД 6: Архитектура ===
    add_content_slide(prs, content_layout,
        "Архитектура матрицы: 5 блоков",
        [
            "Блок 1. Общая информация — название, дата, город, приоритет",
            "Блок 2. Исходные данные — участники, встречи, бюджет",
            "Блок 3. Расчётные метрики — OTS, GRP, CPL, ROI, NPS",
            "Блок 4. Индекс опций — оценка 8 параметров (k=1...5)",
            "Блок 5. Комплексный балл (КБМ) — итоговый рейтинг"
        ])
    
    # === СЛАЙД 7: Метрики ===
    add_content_slide(prs, content_layout,
        "Ключевые метрики",
        [
            "• OTS — охват (сколько могли увидеть бренд)",
            "• GRP — доля с прямым контактом (стенд, сессия)",
            "• CPL — стоимость лида (₽)",
            "• CPA — стоимость целевого действия",
            "• ROI — возврат инвестиций",
            "• NPS — индекс лояльности",
            "• CR — конверсия лид → сделка"
        ])
    
    # === СЛАЙД 8: Индекс опций ===
    add_content_slide(prs, content_layout,
        "Индекс спонсорских опций (Idx)",
        [
            "8 параметров (p1–p8), каждый оценивается k от 1 до 5:",
            "",
            "p1: Выставочный стенд          p5: Сувениры / мерч",
            "p2: Сессия в программе          p6: Активации / игры",
            "p3: Размещение на сайте         p7: Нетворкинг",
            "p4: Логотип на баннерах         p8: Прочее",
            "",
            "Idx = Σ(k) / n — средний балл по всем опциям"
        ])
    
    # === СЛАЙД 9: Матрица k ===
    add_content_slide(prs, content_layout,
        "Матрица коэффициентов (k)",
        [
            "Пример: Выставочный стенд",
            "",
            "k=5  Центральный, у входа, ≥70 кв.м.",
            "k=4  Ключевой павильон, 1-я линия",
            "k=3  2-я линия, видимый",
            "k=2  Угловой, малая площадь",
            "k=1  Задворки, плохая проходимость",
            "",
            "Аналогичная шкала для всех 8 опций"
        ])
    
    # === СЛАЙД 10: КБМ ===
    add_content_slide(prs, content_layout,
        "Комплексный балл мероприятия (КБМ)",
        [
            "КБМ = w1×GRP + w2×Idx + w3×(1/CPL) + w4×NPS + w5×ROI",
            "",
            "Рекомендуемые веса:",
            "• GRP (w1) = 25%  — охват",
            "• Idx (w2) = 20%  — качество пакета",
            "• 1/CPL (w3) = 30%  — эффективность лидов",
            "• NPS (w4) = 15%  — лояльность",
            "• ROI (w5) = 10%  — финансы",
            "",
            "Веса настраиваются под приоритеты компании"
        ])
    
    # === СЛАЙД 11: Разделитель "Пример" ===
    slide11 = prs.slides.add_slide(divider_blue)
    for shape in slide11.placeholders:
        shape.text = "Пример расчёта"
    
    # === СЛАЙД 12: Пример ===
    add_content_slide(prs, content_layout,
        "Сравнение мероприятий 2024-2025",
        [
            "                    ПМГФ      Smart Mining   Белые ночи",
            "",
            "OTS (охват)        34 400          554           375",
            "Касания               186           84            76",
            "GRP                 0,54%       15,16%        20,27%",
            "Idx                  2,75         3,33          2,75",
            "",
            "→ Smart Mining и Белые ночи: выше вовлечённость",
            "   при меньшем охвате"
        ])
    
    # === СЛАЙД 13: Выводы ===
    add_content_slide(prs, content_layout,
        "Выводы и следующие шаги",
        [
            "Что даёт методика:",
            "• Объективное сравнение любых мероприятий",
            "• Обоснование бюджетов на спонсорство",
            "• Выявление эффективных форматов",
            "",
            "Следующие шаги:",
            "• Утвердить веса KPI",
            "• Внедрить для всех мероприятий 2025",
            "• Сформировать рейтинг по итогам года"
        ])
    
    # === СЛАЙД 14: Спасибо ===
    slide14 = prs.slides.add_slide(layouts['1_Спасибо за внимание'])
    
    prs.save(OUTPUT_PATH)
    print(f"✅ Сохранено: {OUTPUT_PATH}")
    print(f"   Слайдов: {len(prs.slides)}")

if __name__ == '__main__':
    create_presentation()

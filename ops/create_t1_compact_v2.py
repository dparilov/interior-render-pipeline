#!/usr/bin/env python3
"""
Компактная версия v2 — текстбоксы вместо placeholder'ов
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

TEMPLATE_PATH = '/home/dima/.openclaw/media/inbound/RUS_Template_T1_окт24_облегченная_версия_2---e96271ef-6531-4673-9d93-5559d2895d14.pptx'
OUTPUT_PATH = '/home/dima/.openclaw/workspace/Metodika_T1_compact_v2.pptx'

def add_title(slide, text):
    """Ищет title placeholder и заполняет"""
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 1:  # TITLE
            shape.text = text
            return
    # Если не нашли — первый placeholder
    if slide.placeholders:
        list(slide.placeholders)[0].text = text

def add_body_text(slide, lines, left=0.7, top=1.5, width=11.9, height=5.5, font_size=16):
    """Добавляет текстовый блок с контентом"""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), 
        Inches(width), Inches(height)
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        
        # Жирный для заголовков секций
        if line.endswith(":") or line.startswith("ПРОБЛЕМА") or line.startswith("РЕШЕНИЕ") or line.startswith("→"):
            p.font.bold = True
        if line.startswith("•"):
            p.level = 0
            p.font.size = Pt(font_size - 1)

def create_presentation():
    prs = Presentation(TEMPLATE_PATH)
    layouts = {layout.name: layout for layout in prs.slide_layouts}
    
    # Удаляем все слайды
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    content_layout = layouts.get('Контентный слайд')
    
    # === СЛАЙД 1: Титульный ===
    slide1 = prs.slides.add_slide(layouts['Титульный слайд 1'])
    for shape in slide1.placeholders:
        if shape.placeholder_format.type == 1:
            shape.text = "Методика оценки эффективности\nспонсорских мероприятий"
        elif shape.placeholder_format.type == 2:
            shape.text = "Управление по работе с бизнес-мероприятиями\nХолдинг Т1 | 2025"
    
    # === СЛАЙД 2: Проблема + Решение ===
    slide2 = prs.slides.add_slide(content_layout)
    add_title(slide2, "Проблема и решение")
    add_body_text(slide2, [
        "ПРОБЛЕМА:",
        "• Нет единой метрики для сравнения мероприятий",
        "• Субъективность оценки спонсорских пакетов",
        "• Сложно обосновать ROI перед руководством",
        "",
        "РЕШЕНИЕ — Матрица оценки эффективности:",
        "• Количественные метрики (охват, лиды, ROI)",
        "• Качественная оценка опций (коэффициенты k)",
        "• Комплексный балл мероприятия (КБМ)",
        "",
        "→ Объективное сравнение любых мероприятий"
    ])
    
    # === СЛАЙД 3: Архитектура ===
    slide3 = prs.slides.add_slide(content_layout)
    add_title(slide3, "Как устроена матрица")
    add_body_text(slide3, [
        "ВВОД ДАННЫХ (ивент-менеджер по данным от организаторов):",
        "• Участники, встречи, контакты, бюджет, доход",
        "• Оценка 8 спонсорских опций (k от 1 до 5)",
        "",
        "                                    ↓",
        "",
        "АВТОМАТИЧЕСКИЙ РАСЧЁТ:",
        "• Метрики: OTS, GRP, CPL, CPA, ROI, NPS",
        "• Индекс опций (Idx) = среднее значение k",
        "• Комплексный балл мероприятия (КБМ)",
        "",
        "                                    ↓",
        "",
        "РЕЗУЛЬТАТ: рейтинг мероприятий для сравнения и планирования"
    ], font_size=15)
    
    # === СЛАЙД 4: КБМ + Пример ===
    slide4 = prs.slides.add_slide(content_layout)
    add_title(slide4, "Комплексный балл и пример расчёта")
    add_body_text(slide4, [
        "ФОРМУЛА КБМ:",
        "КБМ = 25%×GRP + 20%×Idx + 30%×(1/CPL) + 25%×CR",
        "",
        "ПРИМЕР — мероприятия 2024-2025:",
        "",
        "                          ПМГФ        Smart Mining    Белые ночи",
        "OTS (охват)            34 400              554              375",
        "GRP                      0,54%          15,16%          20,27%",
        "Idx                        2,75             3,33             2,75",
        "",
        "→ Smart Mining и Белые ночи: выше вовлечённость (GRP)",
        "   при меньшем охвате — эффективнее для B2B-лидов"
    ], font_size=14)
    
    # === СЛАЙД 5: Выводы ===
    slide5 = prs.slides.add_slide(content_layout)
    add_title(slide5, "Выводы и следующие шаги")
    add_body_text(slide5, [
        "ЧТО ДАЁТ МЕТОДИКА:",
        "• Объективное сравнение любых мероприятий",
        "• Обоснование бюджетов на спонсорство",
        "• База для решений «участвовать / не участвовать»",
        "",
        "ПРОЦЕСС:",
        "• Ивент-менеджер заполняет матрицу по данным от организаторов",
        "• Метрики и КБМ рассчитываются автоматически",
        "• Формируется рейтинг для планирования",
        "",
        "СЛЕДУЮЩИЕ ШАГИ:",
        "• Утвердить веса KPI под цели Т1",
        "• Пилот на мероприятиях Q2 2025",
        "• Рейтинг по итогам года"
    ])
    
    # === СЛАЙД 6: Спасибо ===
    slide6 = prs.slides.add_slide(layouts['1_Спасибо за внимание'])
    
    prs.save(OUTPUT_PATH)
    print(f"✅ Сохранено: {OUTPUT_PATH}")
    print(f"   Слайдов: {len(prs.slides)}")

if __name__ == '__main__':
    create_presentation()

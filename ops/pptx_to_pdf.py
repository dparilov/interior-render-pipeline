#!/usr/bin/env python3
"""
Конвертация PPTX в PDF через reportlab (без LibreOffice)
Читает текст из PPTX и рендерит в PDF
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.colors import HexColor
import os

# Регистрируем шрифт с поддержкой кириллицы
FONT_PATHS = [
    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
    '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
    '/usr/share/fonts/truetype/freefont/FreeSans.ttf',
]

def register_font():
    for path in FONT_PATHS:
        if os.path.exists(path):
            try:
                pdfmetrics.registerFont(TTFont('CustomFont', path))
                return 'CustomFont'
            except:
                continue
    return 'Helvetica'

def pptx_to_pdf(pptx_path, pdf_path):
    prs = Presentation(pptx_path)
    
    # Размер страницы (landscape A4)
    page_width, page_height = landscape(A4)
    
    c = canvas.Canvas(pdf_path, pagesize=(page_width, page_height))
    font_name = register_font()
    
    for slide_num, slide in enumerate(prs.slides, 1):
        # Собираем весь текст со слайда
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text.strip())
        
        # Рисуем номер слайда
        c.setFont(font_name, 10)
        c.setFillColor(HexColor('#888888'))
        c.drawRightString(page_width - 30, 30, f"{slide_num} / {len(prs.slides)}")
        
        # Рисуем текст
        y_position = page_height - 60
        
        for i, text in enumerate(texts):
            lines = text.split('\n')
            for line in lines:
                # Первый текст — заголовок (крупнее)
                if i == 0 and y_position > page_height - 100:
                    c.setFont(font_name, 24)
                    c.setFillColor(HexColor('#0078D7'))
                else:
                    c.setFont(font_name, 12)
                    c.setFillColor(HexColor('#1E1E1E'))
                
                # Обрезаем длинные строки
                if len(line) > 100:
                    line = line[:97] + '...'
                
                c.drawString(40, y_position, line)
                y_position -= 18
                
                if y_position < 50:
                    break
            
            y_position -= 10  # Отступ между блоками
            if y_position < 50:
                break
        
        c.showPage()
    
    c.save()
    print(f"✅ PDF создан: {pdf_path}")

if __name__ == '__main__':
    import sys
    if len(sys.argv) >= 3:
        pptx_to_pdf(sys.argv[1], sys.argv[2])
    else:
        pptx_to_pdf(
            '/home/dima/.openclaw/workspace/Metodika_clean.pptx',
            '/home/dima/.openclaw/workspace/Metodika_clean.pdf'
        )
